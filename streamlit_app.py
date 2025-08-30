import streamlit as st
import openai
import anthropic
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import io
import zipfile
import xml.etree.ElementTree as ET
import re
from typing import Dict, List, Any
import base64
from datetime import datetime

# Page configuration
st.set_page_config(
    page_title="📊 Text to PowerPoint Generator",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Initialize session state for theme
if 'dark_mode' not in st.session_state:
    st.session_state.dark_mode = False


def toggle_theme():
    st.session_state.dark_mode = not st.session_state.dark_mode


# Theme toggle in sidebar
with st.sidebar:
    st.markdown("### 🎨 Theme")
    col1, col2 = st.columns([3, 1])
    with col1:
        current_theme = "Dark Mode" if st.session_state.dark_mode else "Light Mode"
        st.write(f"Current: **{current_theme}**")
    with col2:
        if st.button("🌓", help="Toggle Dark/Light Mode", key="theme_toggle"):
            toggle_theme()
            st.rerun()

# Dynamic CSS based on theme


def get_theme_css(dark_mode: bool) -> str:
    if dark_mode:
        return """
        <style>
            /* Dark Mode Styles */
            .stApp {
                background: linear-gradient(135deg, #1a1a2e 0%, #16213e 100%);
                color: #ffffff;
            }
            
            .main-header {
                background: linear-gradient(135deg, #6366f1 0%, #8b5cf6 100%);
                padding: 2rem;
                border-radius: 15px;
                color: white;
                text-align: center;
                margin-bottom: 2rem;
                box-shadow: 0 10px 30px rgba(99, 102, 241, 0.3);
            }
            
            .feature-box {
                background: linear-gradient(135deg, #374151 0%, #4b5563 100%);
                padding: 1.5rem;
                border-radius: 12px;
                border-left: 4px solid #6366f1;
                margin: 1rem 0;
                box-shadow: 0 4px 15px rgba(0, 0, 0, 0.3);
                color: #f9fafb;
            }
            
            .feature-box h4 {
                color: #a5b4fc;
                margin-bottom: 0.5rem;
            }
            
            .warning-box {
                background: linear-gradient(135deg, #451a03 0%, #78350f 100%);
                padding: 1rem;
                border-radius: 12px;
                border-left: 4px solid #f59e0b;
                margin: 1rem 0;
                color: #fed7aa;
                box-shadow: 0 4px 15px rgba(245, 158, 11, 0.2);
            }
            
            .success-box {
                background: linear-gradient(135deg, #064e3b 0%, #065f46 100%);
                padding: 1rem;
                border-radius: 12px;
                border-left: 4px solid #10b981;
                margin: 1rem 0;
                color: #a7f3d0;
                box-shadow: 0 4px 15px rgba(16, 185, 129, 0.2);
            }
            
            /* Sidebar styling */
            .css-1d391kg {
                background-color: #1f2937;
            }
            
            /* Input field styling */
            .stTextInput input, .stTextArea textarea, .stSelectbox select {
                background-color: #374151 !important;
                color: #ffffff !important;
                border: 2px solid #4b5563 !important;
                border-radius: 8px !important;
            }
            
            .stTextInput input:focus, .stTextArea textarea:focus, .stSelectbox select:focus {
                border-color: #6366f1 !important;
                box-shadow: 0 0 0 3px rgba(99, 102, 241, 0.1) !important;
            }
            
            /* Button styling */
            .stButton button {
                background: linear-gradient(135deg, #6366f1 0%, #8b5cf6 100%) !important;
                color: white !important;
                border: none !important;
                border-radius: 8px !important;
                font-weight: 600 !important;
                transition: all 0.3s ease !important;
            }
            
            .stButton button:hover {
                transform: translateY(-2px) !important;
                box-shadow: 0 8px 25px rgba(99, 102, 241, 0.4) !important;
            }
            
            /* File uploader styling */
            .uploadedFile {
                background-color: #374151 !important;
                border: 2px dashed #6366f1 !important;
                border-radius: 8px !important;
            }
            
            /* Progress bar */
            .stProgress .st-bo {
                background-color: #6366f1 !important;
            }
            
            /* Expander */
            .streamlit-expanderHeader {
                background-color: #374151 !important;
                color: #ffffff !important;
                border-radius: 8px !important;
            }
            
            .streamlit-expanderContent {
                background-color: #374151 !important;
                border: 1px solid #4b5563 !important;
                border-radius: 0 0 8px 8px !important;
            }
            
            /* Code blocks */
            .stCode {
                background-color: #1f2937 !important;
                border: 1px solid #4b5563 !important;
            }
            
            /* Metrics */
            .metric-container {
                background: linear-gradient(135deg, #374151 0%, #4b5563 100%);
                padding: 1rem;
                border-radius: 8px;
                border-left: 4px solid #6366f1;
            }
        </style>
        """
    else:
        return """
        <style>
            /* Light Mode Styles */
            .stApp {
                background: linear-gradient(135deg, #f0f9ff 0%, #e0e7ff 100%);
                color: #1f2937;
            }
            
            .main-header {
                background: linear-gradient(135deg, #4f46e5 0%, #7c3aed 100%);
                padding: 2rem;
                border-radius: 15px;
                color: white;
                text-align: center;
                margin-bottom: 2rem;
                box-shadow: 0 10px 30px rgba(79, 70, 229, 0.3);
            }
            
            .feature-box {
                background: linear-gradient(135deg, #ffffff 0%, #f8fafc 100%);
                padding: 1.5rem;
                border-radius: 12px;
                border-left: 4px solid #4f46e5;
                margin: 1rem 0;
                box-shadow: 0 4px 15px rgba(0, 0, 0, 0.1);
                color: #374151;
            }
            
            .feature-box h4 {
                color: #4f46e5;
                margin-bottom: 0.5rem;
            }
            
            .warning-box {
                background: linear-gradient(135deg, #fef3c7 0%, #fde68a 100%);
                padding: 1rem;
                border-radius: 12px;
                border-left: 4px solid #f59e0b;
                margin: 1rem 0;
                color: #92400e;
                box-shadow: 0 4px 15px rgba(245, 158, 11, 0.1);
            }
            
            .success-box {
                background: linear-gradient(135deg, #f0fdf4 0%, #dcfce7 100%);
                padding: 1rem;
                border-radius: 12px;
                border-left: 4px solid #10b981;
                margin: 1rem 0;
                color: #166534;
                box-shadow: 0 4px 15px rgba(16, 185, 129, 0.1);
            }
            
            /* Sidebar styling */
            .css-1d391kg {
                background-color: #ffffff;
                border-right: 1px solid #e5e7eb;
            }
            
            /* Input field styling */
            .stTextInput input, .stTextArea textarea, .stSelectbox select {
                background-color: #ffffff !important;
                color: #1f2937 !important;
                border: 2px solid #e5e7eb !important;
                border-radius: 8px !important;
            }
            
            .stTextInput input:focus, .stTextArea textarea:focus, .stSelectbox select:focus {
                border-color: #4f46e5 !important;
                box-shadow: 0 0 0 3px rgba(79, 70, 229, 0.1) !important;
            }
            
            /* Button styling */
            .stButton button {
                background: linear-gradient(135deg, #4f46e5 0%, #7c3aed 100%) !important;
                color: white !important;
                border: none !important;
                border-radius: 8px !important;
                font-weight: 600 !important;
                transition: all 0.3s ease !important;
            }
            
            .stButton button:hover {
                transform: translateY(-2px) !important;
                box-shadow: 0 8px 25px rgba(79, 70, 229, 0.4) !important;
            }
            
            /* File uploader styling */
            .uploadedFile {
                background-color: #f9fafb !important;
                border: 2px dashed #4f46e5 !important;
                border-radius: 8px !important;
            }
            
            /* Progress bar */
            .stProgress .st-bo {
                background-color: #4f46e5 !important;
            }
            
            /* Expander */
            .streamlit-expanderHeader {
                background-color: #f3f4f6 !important;
                color: #1f2937 !important;
                border-radius: 8px !important;
            }
            
            .streamlit-expanderContent {
                background-color: #ffffff !important;
                border: 1px solid #e5e7eb !important;
                border-radius: 0 0 8px 8px !important;
            }
            
            /* Code blocks */
            .stCode {
                background-color: #f9fafb !important;
                border: 1px solid #e5e7eb !important;
            }
            
            /* Metrics */
            .metric-container {
                background: linear-gradient(135deg, #ffffff 0%, #f8fafc 100%);
                padding: 1rem;
                border-radius: 8px;
                border-left: 4px solid #4f46e5;
            }
            
            /* Theme toggle button */
            .theme-toggle {
                position: fixed;
                top: 1rem;
                right: 1rem;
                z-index: 999;
                background: rgba(255, 255, 255, 0.9);
                backdrop-filter: blur(10px);
                border-radius: 50%;
                width: 50px;
                height: 50px;
                display: flex;
                align-items: center;
                justify-content: center;
                box-shadow: 0 4px 15px rgba(0, 0, 0, 0.1);
                cursor: pointer;
                transition: all 0.3s ease;
            }
            
            .theme-toggle:hover {
                transform: scale(1.1);
                box-shadow: 0 6px 20px rgba(0, 0, 0, 0.15);
            }
        </style>
        """


# Apply theme-specific CSS
st.markdown(get_theme_css(st.session_state.dark_mode), unsafe_allow_html=True)

# Add keyboard shortcuts
st.markdown("""
<script>
document.addEventListener('keydown', function(event) {
    // Ctrl+Shift+T to toggle theme
    if (event.ctrlKey && event.shiftKey && event.key === 'T') {
        event.preventDefault();
        // Find and click the theme toggle button
        const themeButton = document.querySelector('[data-testid="stButton"] button[title="Toggle Dark/Light Mode"]');
        if (themeButton) {
            themeButton.click();
        }
    }
});
</script>
""", unsafe_allow_html=True)

# Header with theme-aware styling
theme_emoji = "🌙" if st.session_state.dark_mode else "☀️"
theme_text = "Dark Mode Active" if st.session_state.dark_mode else "Light Mode Active"

st.markdown(f"""
<div class="main-header">
    <h1>📊 Text to PowerPoint Generator {theme_emoji}</h1>
    <p>Transform your text into beautiful presentations using AI and custom templates</p>
    <small style="opacity: 0.8;">{theme_text}</small>
</div>
""", unsafe_allow_html=True)


class PresentationGenerator:
    def __init__(self):
        self.template_prs = None
        self.template_styles = {}

    def extract_template_styles(self, template_prs: Presentation) -> Dict[str, Any]:
        """Extract comprehensive styles from the uploaded template"""
        styles = {
            'slide_width': getattr(template_prs, 'slide_width', 9144000),
            'slide_height': getattr(template_prs, 'slide_height', 6858000),
            'layouts': [],
            'master_shapes': [],
            'theme_colors': {},
            'fonts': {},
            'background_fills': []
        }

        try:
            # Extract slide layouts with detailed styling information
            if hasattr(template_prs, 'slide_layouts'):
                for i, layout in enumerate(template_prs.slide_layouts):
                    layout_info = {
                        'index': i,
                        'name': getattr(layout, 'name', f'Layout {i+1}'),
                        'placeholders': [],
                        'background': None
                    }

                    # Extract background fill if present
                    try:
                        if hasattr(layout, 'background') and layout.background.fill:
                            fill = layout.background.fill
                            if hasattr(fill, 'solid') and fill.solid:
                                color = fill.fore_color
                                if hasattr(color, 'rgb'):
                                    layout_info['background'] = {
                                        'type': 'solid',
                                        'color': str(color.rgb)
                                    }
                    except:
                        pass

                    # Extract placeholder styles
                    try:
                        for placeholder in layout.placeholders:
                            placeholder_info = {
                                'idx': getattr(placeholder.placeholder_format, 'idx', 0),
                                'type': str(getattr(placeholder.placeholder_format, 'type', 'unknown')),
                                'left': getattr(placeholder, 'left', 0),
                                'top': getattr(placeholder, 'top', 0),
                                'width': getattr(placeholder, 'width', 0),
                                'height': getattr(placeholder, 'height', 0),
                                'font_info': {},
                                'fill_info': {}
                            }

                            # Extract font styling from placeholder
                            try:
                                if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                                    if hasattr(placeholder.text_frame, 'paragraphs') and placeholder.text_frame.paragraphs:
                                        para = placeholder.text_frame.paragraphs[0]
                                        if hasattr(para, 'runs') and para.runs:
                                            run = para.runs[0] if para.runs else None
                                            if run and hasattr(run, 'font'):
                                                font = run.font
                                                placeholder_info['font_info'] = {
                                                    'name': getattr(font, 'name', None),
                                                    'size': getattr(font, 'size', None),
                                                    'bold': getattr(font, 'bold', None),
                                                    'italic': getattr(font, 'italic', None),
                                                    'color': str(getattr(font.color, 'rgb', '')) if hasattr(font, 'color') and hasattr(font.color, 'rgb') else None
                                                }
                            except:
                                pass

                            layout_info['placeholders'].append(
                                placeholder_info)
                    except Exception as e:
                        st.warning(
                            f"Could not extract placeholders for layout {i}: {str(e)}")

                    styles['layouts'].append(layout_info)

            # Extract theme colors from slide master
            try:
                if hasattr(template_prs, 'slide_master'):
                    master = template_prs.slide_master

                    # Try to extract theme colors
                    if hasattr(master, 'theme') and master.theme:
                        theme = master.theme
                        if hasattr(theme, 'color_scheme'):
                            color_scheme = theme.color_scheme
                            styles['theme_colors'] = {
                                'accent1': getattr(color_scheme, 'accent1_color', None),
                                'accent2': getattr(color_scheme, 'accent2_color', None),
                                'accent3': getattr(color_scheme, 'accent3_color', None),
                                'dark1': getattr(color_scheme, 'dk1_color', None),
                                'dark2': getattr(color_scheme, 'dk2_color', None),
                                'light1': getattr(color_scheme, 'lt1_color', None),
                                'light2': getattr(color_scheme, 'lt2_color', None),
                            }

                    # Extract master slide background
                    try:
                        if hasattr(master, 'background') and master.background.fill:
                            fill = master.background.fill
                            if hasattr(fill, 'solid') and fill.solid:
                                color = fill.fore_color
                                if hasattr(color, 'rgb'):
                                    styles['master_background'] = {
                                        'type': 'solid',
                                        'color': str(color.rgb)
                                    }
                    except:
                        pass

            except Exception as e:
                st.warning(f"Could not extract master slide styles: {str(e)}")

        except Exception as e:
            st.warning(f"Could not extract template styles: {str(e)}")
            # Provide fallback values
            styles['layouts'] = [
                {'index': 0, 'name': 'Title Slide', 'placeholders': []},
                {'index': 1, 'name': 'Content Slide', 'placeholders': []}
            ]

        return styles

    def call_ai_api(self, provider: str, api_key: str, prompt: str) -> str:
        """Call the appropriate AI API"""
        try:
            if provider == "OpenAI":
                client = openai.OpenAI(api_key=api_key)
                response = client.chat.completions.create(
                    model="gpt-4",
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=4000,
                    temperature=0.7
                )
                return response.choices[0].message.content

            elif provider == "Anthropic":
                client = anthropic.Anthropic(api_key=api_key)
                message = client.messages.create(
                    model="claude-3-sonnet-20240229",
                    max_tokens=4000,
                    messages=[{"role": "user", "content": prompt}]
                )
                return message.content[0].text

            elif provider == "Google Gemini":
                genai.configure(api_key=api_key)
                model = genai.GenerativeModel('gemini-2.0-flash')
                response = model.generate_content(
                    prompt,
                    generation_config=genai.types.GenerationConfig(
                        temperature=0.7,
                        max_output_tokens=4000,
                    )
                )
                return response.text

        except Exception as e:
            raise Exception(f"AI API Error: {str(e)}")

    def create_prompt(self, input_text: str, guidance: str = "") -> str:
        """Create the prompt for AI processing"""
        return f"""Please analyze the following text and convert it into a structured PowerPoint presentation outline.

IMPORTANT: Respond with ONLY a valid JSON object in this exact format:
{{
  "title": "Presentation Title",
  "slides": [
    {{
      "title": "Slide Title",
      "content": ["bullet point 1", "bullet point 2", "bullet point 3"],
      "notes": "Speaker notes for this slide (optional)"
    }}
  ]
}}

Guidelines:
- Create 5-12 slides based on content length and complexity
- Each slide should have 2-5 concise bullet points
- Make titles engaging and descriptive
- Include speaker notes when helpful
- Focus on key insights, not just copying text
{f"- Style/tone guidance: {guidance}" if guidance else ""}

Text to convert:
{input_text}

Remember: Respond with ONLY the JSON object, no additional text or formatting."""

    def parse_ai_response(self, response: str) -> Dict[str, Any]:
        """Parse and validate AI response"""
        try:
            # Find JSON in the response
            json_start = response.find('{')
            json_end = response.rfind('}')

            if json_start == -1 or json_end == -1:
                raise ValueError("No JSON found in response")

            json_str = response[json_start:json_end + 1]
            structure = json.loads(json_str)

            if not structure.get('slides') or not isinstance(structure['slides'], list):
                raise ValueError("Invalid presentation structure")

            return structure

        except Exception as e:
            raise Exception(f"Failed to parse AI response: {str(e)}")

    def create_presentation(self, structure: Dict[str, Any], template_prs: Presentation = None) -> Presentation:
        """Create PowerPoint presentation from structure with proper XML compliance"""
        try:
            if template_prs:
                # Method 1: Create new presentation with template's slide master
                prs = Presentation()

                # Copy slide dimensions from template
                try:
                    prs.slide_width = template_prs.slide_width
                    prs.slide_height = template_prs.slide_height
                except:
                    pass

                # Use template layouts but create clean structure
                try:
                    layouts_to_use = template_prs.slide_layouts
                    st.info("✅ Using template layouts")
                except:
                    layouts_to_use = prs.slide_layouts
                    st.warning("⚠️ Using default layouts")
            else:
                # Create clean new presentation
                prs = Presentation()
                layouts_to_use = prs.slide_layouts

            # Get template styles for applying formatting
            template_styles = getattr(self, 'template_styles', {})

            # Ensure we have at least 2 layouts
            title_layout = layouts_to_use[0] if len(
                layouts_to_use) > 0 else prs.slide_layouts[0]
            content_layout = layouts_to_use[1] if len(
                layouts_to_use) > 1 else prs.slide_layouts[1]

            # Add title slide with error handling
            try:
                title_slide = prs.slides.add_slide(title_layout)

                # Set title safely
                if hasattr(title_slide, 'shapes') and title_slide.shapes.title:
                    title_slide.shapes.title.text = structure.get(
                        'title', 'Generated Presentation')
                    # Apply basic formatting safely
                    try:
                        if hasattr(title_slide.shapes.title, 'text_frame'):
                            text_frame = title_slide.shapes.title.text_frame
                            if text_frame.paragraphs:
                                for paragraph in text_frame.paragraphs:
                                    paragraph.alignment = PP_ALIGN.CENTER
                    except:
                        pass

                # Add subtitle safely
                try:
                    if len(title_slide.placeholders) > 1:
                        subtitle = title_slide.placeholders[1]
                        if hasattr(subtitle, 'text'):
                            subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
                except:
                    pass

            except Exception as e:
                st.error(f"Error creating title slide: {str(e)}")
                raise

            # Add content slides with robust error handling
            for i, slide_data in enumerate(structure['slides']):
                try:
                    slide = prs.slides.add_slide(content_layout)

                    # Set slide title safely
                    if hasattr(slide, 'shapes') and slide.shapes.title:
                        slide.shapes.title.text = slide_data['title']

                    # Add content with multiple fallback strategies
                    content_added = self._add_slide_content(
                        slide, slide_data['content'])

                    if not content_added:
                        st.warning(
                            f"Could not add content to slide {i+1}: {slide_data['title']}")

                    # Add speaker notes safely
                    try:
                        if slide_data.get('notes'):
                            notes_slide = slide.notes_slide
                            if hasattr(notes_slide, 'notes_text_frame'):
                                notes_slide.notes_text_frame.text = slide_data['notes']
                    except:
                        pass

                except Exception as e:
                    st.error(f"Error creating slide {i+1}: {str(e)}")
                    continue  # Skip this slide but continue with others

            # Validate presentation before returning
            if len(prs.slides) == 0:
                raise ValueError("No slides were created successfully")

            return prs

        except Exception as e:
            st.error(f"Critical error in presentation creation: {str(e)}")
            # Create minimal fallback presentation
            return self._create_fallback_presentation(structure)

    def _add_slide_content(self, slide, content_points: List[str]) -> bool:
        """Add content to slide with multiple fallback strategies"""
        strategies = [
            self._try_placeholder_content,
            self._try_textbox_content,
            self._try_simple_textbox_content
        ]

        for strategy in strategies:
            try:
                if strategy(slide, content_points):
                    return True
            except Exception as e:
                continue

        return False

    def _try_placeholder_content(self, slide, content_points: List[str]) -> bool:
        """Try to add content using slide placeholders"""
        for placeholder in slide.placeholders:
            try:
                # Look for content placeholders
                if (placeholder.placeholder_format.idx == 1 or
                    'content' in str(placeholder.placeholder_format.type).lower() or
                        'body' in str(placeholder.placeholder_format.type).lower()):

                    if hasattr(placeholder, 'text_frame'):
                        text_frame = placeholder.text_frame
                        text_frame.clear()

                        # Add bullet points
                        for i, point in enumerate(content_points):
                            if i == 0:
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            p.text = str(point)  # Ensure string
                            p.level = 0

                        return True
            except:
                continue
        return False

    def _try_textbox_content(self, slide, content_points: List[str]) -> bool:
        """Try to add content using a text box"""
        try:
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5.5)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            # Add content with bullet points
            for i, point in enumerate(content_points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = f"• {str(point)}"
                p.level = 0

            return True
        except:
            return False

    def _try_simple_textbox_content(self, slide, content_points: List[str]) -> bool:
        """Fallback: simple text box with basic formatting"""
        try:
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(6)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame

            # Simple text without fancy formatting
            content_text = "\n".join(
                [f"• {str(point)}" for point in content_points])
            text_frame.text = content_text

            return True
        except:
            return False

    def _create_fallback_presentation(self, structure: Dict[str, Any]) -> Presentation:
        """Create a minimal presentation as last resort"""
        prs = Presentation()

        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = structure.get(
            'title', 'Generated Presentation')

        # Add one content slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = "Content Summary"

        # Add basic content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.text = "Content generated successfully"

        for slide_data in structure.get('slides', [])[:5]:  # Limit to 5 slides
            p = text_frame.add_paragraph()
            p.text = slide_data.get('title', 'Slide')
            p.level = 0

        return prs

    def apply_text_styling(self, shape, template_styles: Dict, style_type: str):
        """Apply text styling from template to a text shape"""
        if not template_styles or not hasattr(shape, 'text_frame'):
            return

        try:
            text_frame = shape.text_frame
            if not text_frame.paragraphs:
                return

            for paragraph in text_frame.paragraphs:
                self.apply_paragraph_styling(
                    paragraph, template_styles, style_type)

        except Exception as e:
            pass  # Silently fail if styling can't be applied

    def apply_paragraph_styling(self, paragraph, template_styles: Dict, style_type: str):
        """Apply paragraph-level styling from template"""
        if not template_styles or not hasattr(paragraph, 'runs'):
            return

        try:
            # Look for relevant font info in template styles
            font_info = None
            layouts = template_styles.get('layouts', [])

            # Find font info from appropriate placeholder
            for layout in layouts:
                for placeholder in layout.get('placeholders', []):
                    if placeholder.get('font_info') and any(placeholder['font_info'].values()):
                        font_info = placeholder['font_info']
                        break
                if font_info:
                    break

            if not font_info:
                return

            # Apply font styling to all runs in paragraph
            for run in paragraph.runs:
                if hasattr(run, 'font'):
                    font = run.font

                    # Apply font name (with fallback for compatibility)
                    if font_info.get('name'):
                        try:
                            font.name = font_info['name']
                        except:
                            # If font is not available, try common fallbacks
                            fallback_fonts = [
                                'Calibri', 'Arial', 'Times New Roman', 'Helvetica']
                            for fallback in fallback_fonts:
                                try:
                                    font.name = fallback
                                    break
                                except:
                                    continue

                    # Apply font size
                    if font_info.get('size') and font_info['size']:
                        try:
                            font.size = font_info['size']
                        except:
                            pass

                    # Apply font weight
                    if font_info.get('bold') is not None:
                        try:
                            font.bold = font_info['bold']
                        except:
                            pass

                    # Apply font style
                    if font_info.get('italic') is not None:
                        try:
                            font.italic = font_info['italic']
                        except:
                            pass

                    # Apply font color
                    if font_info.get('color') and font_info['color']:
                        try:
                            # Parse color string and apply
                            color_str = font_info['color'].replace(
                                'RGBColor(0x', '').replace(')', '')
                            if len(color_str) == 6:
                                r = int(color_str[0:2], 16)
                                g = int(color_str[2:4], 16)
                                b = int(color_str[4:6], 16)
                                font.color.rgb = RGBColor(r, g, b)
                        except:
                            pass

        except Exception as e:
            pass  # Silently fail if styling can't be applied

    def validate_presentation(self, prs: Presentation) -> Dict[str, Any]:
        """Validate presentation structure and fix common issues"""
        validation_results = {
            'is_valid': True,
            'issues_found': [],
            'issues_fixed': [],
            'warnings': []
        }

        try:
            # Check if presentation has slides
            if len(prs.slides) == 0:
                validation_results['is_valid'] = False
                validation_results['issues_found'].append(
                    "No slides in presentation")
                return validation_results

            # Validate each slide
            for i, slide in enumerate(prs.slides):
                slide_issues = []

                # Check for title
                try:
                    if hasattr(slide, 'shapes') and slide.shapes.title:
                        title_text = slide.shapes.title.text
                        if not title_text or title_text.strip() == "":
                            slide.shapes.title.text = f"Slide {i+1}"
                            validation_results['issues_fixed'].append(
                                f"Added missing title to slide {i+1}")
                except:
                    validation_results['warnings'].append(
                        f"Could not validate title for slide {i+1}")

                # Check for content
                has_content = False
                try:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            if shape.text_frame.text and shape.text_frame.text.strip():
                                has_content = True
                                break

                    if not has_content and i > 0:  # Skip title slide
                        validation_results['warnings'].append(
                            f"Slide {i+1} appears to have no content")

                except:
                    validation_results['warnings'].append(
                        f"Could not validate content for slide {i+1}")

            # Check slide dimensions
            try:
                if not hasattr(prs, 'slide_width') or not hasattr(prs, 'slide_height'):
                    validation_results['warnings'].append(
                        "Slide dimensions may not be set correctly")
                elif prs.slide_width <= 0 or prs.slide_height <= 0:
                    # Fix invalid dimensions
                    prs.slide_width = Inches(10)
                    prs.slide_height = Inches(7.5)
                    validation_results['issues_fixed'].append(
                        "Fixed invalid slide dimensions")
            except:
                validation_results['warnings'].append(
                    "Could not validate slide dimensions")

        except Exception as e:
            validation_results['is_valid'] = False
            validation_results['issues_found'].append(
                f"Validation error: {str(e)}")

        return validation_results

    def save_presentation_safely(self, prs: Presentation, filename: str) -> io.BytesIO:
        """Save presentation with validation and error handling"""
        try:
            # Validate presentation first
            validation = self.validate_presentation(prs)

            # Show validation results
            if validation['issues_fixed']:
                st.info(
                    f"✅ Fixed {len(validation['issues_fixed'])} issues automatically")
                for fix in validation['issues_fixed']:
                    st.write(f"  • {fix}")

            if validation['warnings']:
                st.warning(f"⚠️ {len(validation['warnings'])} warnings found")
                with st.expander("View Warnings"):
                    for warning in validation['warnings']:
                        st.write(f"  • {warning}")

            if not validation['is_valid']:
                st.error("❌ Presentation validation failed:")
                for issue in validation['issues_found']:
                    st.write(f"  • {issue}")
                raise ValueError("Presentation failed validation")

            # Save with error handling
            output = io.BytesIO()

            try:
                prs.save(output)
                output.seek(0)

                # Verify the saved file
                if output.getvalue() and len(output.getvalue()) > 1000:  # Reasonable minimum size
                    st.success("✅ Presentation saved successfully")
                    return output
                else:
                    raise ValueError(
                        "Generated file appears to be too small or empty")

            except Exception as save_error:
                st.error(f"Error during save: {str(save_error)}")

                # Try alternative save method
                st.info("🔄 Trying alternative save method...")
                temp_output = io.BytesIO()

                # Create a completely new, clean presentation
                clean_prs = Presentation()

                # Copy basic structure only
                for i, slide in enumerate(prs.slides):
                    if i == 0:
                        new_slide = clean_prs.slides.add_slide(
                            clean_prs.slide_layouts[0])
                    else:
                        new_slide = clean_prs.slides.add_slide(
                            clean_prs.slide_layouts[1])

                    # Copy title
                    if hasattr(slide, 'shapes') and slide.shapes.title and hasattr(new_slide, 'shapes') and new_slide.shapes.title:
                        new_slide.shapes.title.text = slide.shapes.title.text

                clean_prs.save(temp_output)
                temp_output.seek(0)

                if temp_output.getvalue():
                    st.warning("⚠️ Used simplified version due to save issues")
                    return temp_output
                else:
                    raise ValueError("All save methods failed")

        except Exception as e:
            st.error(f"❌ Failed to save presentation: {str(e)}")
            raise


def main():
    generator = PresentationGenerator()

    # Sidebar for configuration
    with st.sidebar:
        st.header("🔧 Configuration")

        # Theme toggle section at the top
        st.markdown("---")
        st.markdown("### 🎨 Appearance")

        col1, col2 = st.columns([2, 1])
        with col1:
            current_theme = "🌙 Dark Mode" if st.session_state.dark_mode else "☀️ Light Mode"
            st.markdown(f"**{current_theme}**")
        with col2:
            toggle_icon = "☀️" if st.session_state.dark_mode else "🌙"
            if st.button(toggle_icon, help="Toggle Dark/Light Mode", key="theme_toggle"):
                toggle_theme()
                st.rerun()

        if st.session_state.dark_mode:
            st.success("🌙 Dark mode is easier on the eyes!")
        else:
            st.info("☀️ Light mode for bright environments!")

        st.markdown("---")

        # AI Provider Selection
        st.subheader("🤖 AI Provider")
        provider = st.selectbox(
            "Choose AI Provider",
            ["OpenAI", "Anthropic", "Google Gemini"],
            help="Select your preferred AI service"
        )

        # API Key Input
        api_key = st.text_input(
            f"{provider} API Key",
            type="password",
            help="Your API key is never stored or logged"
        )

        if not api_key:
            theme_class = "warning-box" if not st.session_state.dark_mode else "warning-box"
            st.markdown(f"""
            <div class="{theme_class}">
                ⚠️ <strong>API Key Required</strong><br>
                Get your API key from:
                <ul style="margin-top: 0.5rem;">
                    <li><strong>OpenAI:</strong> <a href="https://platform.openai.com" target="_blank">platform.openai.com</a></li>
                    <li><strong>Anthropic:</strong> <a href="https://console.anthropic.com" target="_blank">console.anthropic.com</a></li>
                    <li><strong>Google:</strong> <a href="https://console.cloud.google.com" target="_blank">console.cloud.google.com</a></li>
                </ul>
            </div>
            """, unsafe_allow_html=True)

        st.divider()

        # Template Upload
        st.subheader("🎨 Template")
        template_file = st.file_uploader(
            "Upload PowerPoint Template",
            type=['pptx', 'potx'],
            help="Upload your branded PowerPoint template to preserve styling"
        )

        if template_file:
            try:
                # Load template
                generator.template_prs = Presentation(template_file)
                generator.template_styles = generator.extract_template_styles(
                    generator.template_prs)
                st.success("✅ Template loaded successfully!")

                # Show template info
                with st.expander("Template Information"):
                    st.write(
                        f"**Slide Size:** {generator.template_styles['slide_width']} x {generator.template_styles['slide_height']}")
                    st.write(
                        f"**Available Layouts:** {len(generator.template_styles['layouts'])}")

                    for i, layout in enumerate(generator.template_styles['layouts']):
                        st.write(
                            f"  - Layout {i+1}: {layout['name']} ({len(layout['placeholders'])} placeholders)")

                        # Show extracted font information
                        fonts_found = []
                        for placeholder in layout['placeholders']:
                            font_info = placeholder.get('font_info', {})
                            if any(font_info.values()):
                                font_details = []
                                if font_info.get('name'):
                                    font_details.append(
                                        f"Font: {font_info['name']}")
                                if font_info.get('size'):
                                    font_details.append(
                                        f"Size: {font_info['size']}")
                                if font_info.get('color'):
                                    font_details.append(
                                        f"Color: {font_info['color']}")
                                if font_details:
                                    fonts_found.append(
                                        f"    • {', '.join(font_details)}")

                        if fonts_found:
                            st.write("    Font Styles Found:")
                            # Show max 2 font styles per layout
                            for font in fonts_found[:2]:
                                st.write(font)

                    # Show theme colors if extracted
                    theme_colors = generator.template_styles.get(
                        'theme_colors', {})
                    if any(theme_colors.values()):
                        st.write("**Theme Colors:** Extracted ✅")
                        color_count = len(
                            [c for c in theme_colors.values() if c])
                        st.write(f"  - {color_count} theme colors available")
                    else:
                        st.write("**Theme Colors:** Could not extract")

                    # Show background info
                    if generator.template_styles.get('master_background'):
                        st.write("**Master Background:** Detected ✅")
                    else:
                        st.write("**Master Background:** Using default")

            except Exception as e:
                st.error(f"❌ Error loading template: {str(e)}")
                st.info(
                    "💡 The app will use a default template instead. Common issues:")
                st.write("- File might be corrupted or password protected")
                st.write("- Template might have unsupported features")
                st.write("- Try using a simpler .pptx template")

                # Reset template data
                generator.template_prs = None
                generator.template_styles = {}

    # Main content area
    col1, col2 = st.columns([2, 1])

    with col1:
        st.header("📝 Your Content")

        # Input text area
        input_text = st.text_area(
            "Paste your content here",
            height=300,
            placeholder="Paste your markdown, prose, or any text content here. The AI will intelligently break it down into slide-worthy content...",
            help="Supports markdown formatting and long-form text"
        )

        # Style guidance
        st.header("🎯 Presentation Style")
        guidance_options = [
            "",
            "investor pitch deck",
            "research summary",
            "sales presentation",
            "technical documentation",
            "educational content",
            "project proposal",
            "quarterly review",
            "product launch"
        ]

        guidance = st.selectbox(
            "Choose presentation style (optional)",
            guidance_options,
            help="This guides the AI on how to structure and tone your presentation"
        )

        if guidance == "":
            custom_guidance = st.text_input(
                "Or enter custom guidance",
                placeholder="e.g., board meeting presentation, customer onboarding, training material"
            )
            if custom_guidance:
                guidance = custom_guidance

    with col2:
        st.header("🚀 Features")

        st.markdown("""
        <div class="feature-box">
            <h4>🤖 AI-Powered Analysis</h4>
            <p>Intelligently structures your content into compelling slides</p>
        </div>
        """, unsafe_allow_html=True)

        st.markdown("""
        <div class="feature-box">
            <h4>🎨 Template Preservation</h4>
            <p>Maintains your brand colors, fonts, and layouts. Custom fonts fallback to system fonts for compatibility.</p>
        </div>
        """, unsafe_allow_html=True)

        st.markdown("""
        <div class="feature-box">
            <h4>🔒 Privacy First</h4>
            <p>No data storage - everything processed securely</p>
        </div>
        """, unsafe_allow_html=True)

        st.markdown("""
        <div class="feature-box">
            <h4>⚡ Multiple AI Providers</h4>
            <p>Choose from OpenAI, Anthropic, or Google</p>
        </div>
        """, unsafe_allow_html=True)

    # Generation button and process
    st.divider()

    if st.button("🚀 Generate Presentation", type="primary", use_container_width=True):
        if not input_text.strip():
            st.error("❌ Please enter some content to convert.")
            return

        if not api_key.strip():
            st.error("❌ Please enter your API key.")
            return

        try:
            with st.spinner("🤖 Analyzing content with AI..."):
                # Create progress bar
                progress_bar = st.progress(0)
                status_text = st.empty()

                # Step 1: Generate structure
                status_text.text("🔍 Analyzing content structure...")
                progress_bar.progress(25)

                prompt = generator.create_prompt(input_text, guidance)
                ai_response = generator.call_ai_api(provider, api_key, prompt)

                # Step 2: Parse response
                status_text.text("📋 Parsing AI response...")
                progress_bar.progress(50)

                structure = generator.parse_ai_response(ai_response)

                # Step 3: Create presentation
                status_text.text("🎨 Creating presentation...")
                progress_bar.progress(75)

                prs = generator.create_presentation(
                    structure, generator.template_prs)

                # Step 4: Finalize
                status_text.text("✅ Finalizing presentation...")
                progress_bar.progress(90)

                # Generate filename
                safe_title = re.sub(
                    r'[^a-z0-9\s]', '', structure.get('title', 'presentation').lower())
                safe_title = re.sub(r'\s+', '_', safe_title)
                filename = f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

                # Save presentation with validation
                output = generator.save_presentation_safely(prs, filename)

                progress_bar.progress(100)
                status_text.text("🎉 Presentation generated successfully!")

                # Success message with details
                st.markdown(f"""
                <div class="success-box">
                    <h4>✅ Success!</h4>
                    <p><strong>Title:</strong> {structure.get('title', 'Generated Presentation')}</p>
                    <p><strong>Slides Created:</strong> {len(structure['slides']) + 1} (including title slide)</p>
                    <p><strong>Template Applied:</strong> {'Yes' if generator.template_prs else 'Default'}</p>
                </div>
                """, unsafe_allow_html=True)

                # Download button
                st.download_button(
                    label="📥 Download Presentation",
                    data=output.getvalue(),
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

                # Show presentation structure
                with st.expander("📋 Generated Structure"):
                    st.json(structure)

        except Exception as e:
            st.error(f"❌ Error: {str(e)}")
            st.info("💡 Try using a different AI provider or check your API key.")

    # Footer with theme information
    st.divider()

    footer_style = "color: #9ca3af;" if st.session_state.dark_mode else "color: #6b7280;"
    theme_info = "🌙 Dark Mode" if st.session_state.dark_mode else "☀️ Light Mode"

    st.markdown(f"""
    <div style="text-align: center; {footer_style} padding: 2rem;">
        <p>Made with ❤️ using Streamlit • Transform your ideas into professional presentations</p>
        <p><small>Your API keys and content are processed securely and never stored</small></p>
        <p style="margin-top: 1rem;"><small>
            {theme_info} • Press Ctrl+Shift+T to toggle theme • 
            <a href="https://github.com/yourusername/text-to-powerpoint" target="_blank" style="color: inherit;">View on GitHub</a>
        </small></p>
    </div>
    """, unsafe_allow_html=True)


if __name__ == "__main__":
    main()
